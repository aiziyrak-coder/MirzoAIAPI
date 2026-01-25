import os
import base64
import mammoth
import docx
from pptx import Presentation
from PyPDF2 import PdfReader
from django.conf import settings
from django.core.files.uploadedfile import InMemoryUploadedFile
from io import BytesIO
import PIL.Image


def process_file_to_parts(file_path, file_name, mime_type):
    """
    Process uploaded file and return parts for Gemini API
    """
    parts = []
    
    try:
        # Handle DOCX (Word)
        if mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
            with open(file_path, 'rb') as f:
                result = mammoth.extract_raw_text(f)
                parts.append({
                    'text': f"[HUJJAT TARKIBI: {file_name}]:\n{result.value}"
                })
            return parts
        
        # Handle PPTX (PowerPoint)
        if mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
            prs = Presentation(file_path)
            ppt_text = ''
            for i, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                if slide_text:
                    ppt_text += f"--- SLIDE {i} ---\n" + "\n".join(slide_text) + "\n\n"
            
            parts.append({
                'text': f"[PREZENTATSIYA TARKIBI: {file_name}]:\n{ppt_text}"
            })
            return parts
        
        # Handle PDF
        if mime_type == 'application/pdf':
            pdf_text = ''
            try:
                with open(file_path, 'rb') as f:
                    reader = PdfReader(f)
                    for page in reader.pages:
                        pdf_text += page.extract_text() + "\n"
            except Exception as e:
                pdf_text = f"[PDF: {file_name}] - PDF o'qib bo'lmadi: {str(e)}"
            
            parts.append({
                'text': f"[PDF TARKIBI: {file_name}]:\n{pdf_text}"
            })
            return parts
        
        # Handle Images (PNG, JPEG, WEBP) - return as inline_data for Gemini
        supported_image_types = ['image/png', 'image/jpeg', 'image/jpg', 'image/webp']
        if mime_type in supported_image_types:
            with open(file_path, 'rb') as f:
                image_data = base64.b64encode(f.read()).decode('utf-8')
                parts.append({
                    'inline_data': {
                        'mime_type': mime_type,
                        'data': image_data
                    }
                })
            return parts
        
        # Handle plain text
        if mime_type.startswith('text/'):
            with open(file_path, 'r', encoding='utf-8') as f:
                text_content = f.read()
                parts.append({
                    'text': f"Qo'shimcha ilova qilingan hujjat ({file_name}):\n{text_content}"
                })
            return parts
        
    except Exception as e:
        print(f"Error processing file {file_name}: {str(e)}")
        parts.append({
            'text': f"[HUJJAT: {file_name}] - O'qishda xatolik: {str(e)}"
        })
    
    return parts


def process_files_to_parts(files):
    """
    Process multiple files and return combined parts
    """
    all_parts = []
    
    for file in files:
        # Save file temporarily if needed
        file_path = None
        
        if hasattr(file, 'temporary_file_path'):
            file_path = file.temporary_file_path()
        else:
            # Save to temp location
            temp_path = os.path.join(settings.UPLOAD_DIR, file.name)
            os.makedirs(os.path.dirname(temp_path), exist_ok=True)
            
            with open(temp_path, 'wb+') as destination:
                for chunk in file.chunks():
                    destination.write(chunk)
            file_path = temp_path
        
        try:
            parts = process_file_to_parts(file_path, file.name, file.content_type or 'application/octet-stream')
            all_parts.extend(parts)
        except Exception as e:
            print(f"Error processing file {file.name}: {str(e)}")
            all_parts.append({
                'text': f"[HUJJAT: {file.name}] - Xatolik: {str(e)}"
            })
        finally:
            # Clean up temp file
            if file_path and os.path.exists(file_path) and file_path.startswith(str(settings.UPLOAD_DIR)):
                try:
                    os.remove(file_path)
                except:
                    pass
    
    return all_parts
